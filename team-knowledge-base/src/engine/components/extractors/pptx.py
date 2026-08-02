from pathlib import Path

from pptx import Presentation

from src.engine.components.extractors.base import BaseExtractor


class PPTXExtractor(BaseExtractor):
    """提取 PowerPoint (.pptx) 文件文本内容。"""

    SUPPORTED_EXTENSIONS = {".pptx"}

    def extract(self, file_path: Path) -> str:
        self._ensure_file_exists(file_path)
        try:
            prs = Presentation(str(file_path))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                texts.append(text)
                if texts:
                    slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(slides_text)
        except Exception as e:
            raise ValueError(f"PPTX 解析失败: {file_path} — {e}") from e
