"""生成测试用噪声 PPT：季度汇报演示文稿（中文）"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = r"D:\knowledge-vault\test-dirty"
os.makedirs(OUTPUT_DIR, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# 颜色
BLUE = RGBColor(0x1A, 0x56, 0xDB)
DARK = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0x99, 0x99, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)


def add_footer(slide, page_num, total=12):
    """每页重复的页脚噪声"""
    # 公司名
    txBox = slide.shapes.add_textbox(Cm(1), Cm(17.5), Cm(8), Cm(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "星辰科技有限公司 | 内部机密"
    p.font.size = Pt(8)
    p.font.color.rgb = GREY

    # 页码
    txBox2 = slide.shapes.add_textbox(Cm(30), Cm(17.5), Cm(3), Cm(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"{page_num} / {total}"
    p2.font.size = Pt(8)
    p2.font.color.rgb = GREY
    p2.alignment = PP_ALIGN.RIGHT

    # 文档编号
    txBox3 = slide.shapes.add_textbox(Cm(12), Cm(17.5), Cm(10), Cm(0.8))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "MKT-2024-Q2-REPORT | V1.3 | 2024-07-05"
    p3.font.size = Pt(7)
    p3.font.color.rgb = GREY
    p3.alignment = PP_ALIGN.CENTER


def add_title_slide(title, subtitle):
    """标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 标题
    txBox = slide.shapes.add_textbox(Cm(3), Cm(5), Cm(28), Cm(4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK
    # 副标题
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = GREY
    return slide


def add_content_slide(title, bullets, notes=""):
    """内容页：标题 + bullet 列表"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 标题
    txBox = slide.shapes.add_textbox(Cm(1.5), Cm(0.8), Cm(30), Cm(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLUE
    # 内容
    txBox2 = slide.shapes.add_textbox(Cm(2), Cm(3.2), Cm(29), Cm(13))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)
    # 备注（speaker notes）
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_image_placeholder_slide(title, caption):
    """图片占位页：只有一个框和标题，文字提取几乎为空"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 标题
    txBox = slide.shapes.add_textbox(Cm(1.5), Cm(0.8), Cm(30), Cm(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLUE
    # 图片占位框
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1, Cm(3), Cm(3.5), Cm(27), Cm(12)  # MSO_SHAPE.RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GREY
    shape.line.color.rgb = GREY
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = caption
    p.font.size = Pt(14)
    p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.CENTER
    return slide


# ============================================================
# Slide 1: 封面
# ============================================================
slide = add_title_slide(
    "2024年Q2市场营销工作汇报",
    "市场营销部 | 汇报人：林小燕 | 2024年7月5日"
)
# 封面额外噪声
txBox = slide.shapes.add_textbox(Cm(3), Cm(12), Cm(28), Cm(4))
tf = txBox.text_frame
cover_noise = [
    "星辰科技有限公司 STAR CHEN TECHNOLOGY CO., LTD.",
    "地址：北京市海淀区中关村科技园区创新大厦18层",
    "本文档包含商业机密信息，未经授权严禁复制、传播或向第三方披露",
    "文档编号：MKT-2024-Q2-REPORT | 版本：V1.3 | 密级：内部机密",
    "© 2024 星辰科技有限公司 版权所有",
]
for i, line in enumerate(cover_noise):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(9)
    p.font.color.rgb = GREY
add_footer(slide, 1)

# ============================================================
# Slide 2: 目录
# ============================================================
slide = add_content_slide("汇报提纲", [
    "一、Q2市场活动概览",
    "二、品牌推广成效",
    "三、线索获取与转化",
    "四、竞品动态分析",
    "五、Q3工作计划",
    "六、预算执行情况",
    "",
    "（各章节详细内容见后续页面）",
])
add_footer(slide, 2)

# ============================================================
# Slide 3: 概览（碎片化 bullet，缺乏完整语句）
# ============================================================
slide = add_content_slide("一、Q2市场活动概览", [
    "• 线下活动：12场",
    "• 线上直播：8场",
    "• 行业展会：3场（CITE 2024 / 数博会 / 世界人工智能大会）",
    "• 白皮书发布：2份",
    "• 媒体曝光：156篇",
    "• 总触达人次：约 85万",
    "",
    "关键结论：",
    "• 品牌知名度 ↑ 12%（vs Q1）",
    "• 获客成本 ↓ 8%",
    "• MQL → SQL 转化率：23%（目标 25%，未达标）",
], notes="备注：MQL未达标主要原因是5月份官网改版导致表单提交异常，已修复。6月转化率恢复到26%。")
add_footer(slide, 3)

# ============================================================
# Slide 4: 图片页（文字几乎为空）
# ============================================================
slide = add_image_placeholder_slide(
    "Q2市场活动全景图",
    "[此处为活动照片拼图 - 12场线下活动精选]\n（图片未能提取）"
)
add_footer(slide, 4)

# ============================================================
# Slide 5: 品牌推广（有数据但碎片化）
# ============================================================
slide = add_content_slide("二、品牌推广成效", [
    "搜索引擎：",
    "  - 品牌词搜索量：月均 12,000 次（+18% QoQ）",
    "  - SEM 点击率：4.2%（行业均值 3.1%）",
    "",
    "社交媒体：",
    "  - 微信公众号：粉丝 8.2万（+6,000）",
    "  - 知乎专栏：阅读量 45万",
    "  - B站技术视频：播放量 120万",
    "",
    "行业影响力：",
    "  - 受邀演讲：5次（含1次 keynote）",
    "  - 行业标准参编：2项",
    "  - 获奖：「年度创新数据平台」（中国信通院）",
], notes="备注：B站播放量主要来自3月份发布的《数据中台从0到1》系列，长尾效应明显。")
add_footer(slide, 5)

# ============================================================
# Slide 6: 又一个图片页
# ============================================================
slide = add_image_placeholder_slide(
    "品牌搜索指数趋势（2024 H1）",
    "[此处为百度指数截图]\n（图表未能提取）"
)
add_footer(slide, 6)

# ============================================================
# Slide 7: 线索转化（表格数据）
# ============================================================
slide = add_content_slide("三、线索获取与转化", [
    "渠道        | 线索数  | MQL   | SQL  | 成交  | 转化率",
    "─────────────────────────────────────────────────",
    "官网表单    | 2,340  | 1,870 | 420  | 85   | 3.6%",
    "线下活动    | 1,560  | 1,200 | 380  | 92   | 5.9%",
    "线上直播    | 3,200  | 2,100 | 350  | 48   | 1.5%",
    "合作伙伴    | 890    | 750   | 280  | 110  | 12.4%",
    "老客户转介  | 420    | 380   | 210  | 95   | 22.6%",
    "─────────────────────────────────────────────────",
    "合计        | 8,410  | 6,300 | 1,640| 430  | 5.1%",
    "",
    "目标完成率：线索量 105% | SQL 94% | 成交额 112%",
], notes="备注：合作伙伴渠道转化率最高但量最少，Q3计划拓展3家新合作伙伴。线上直播量大但质量低，考虑调整直播策略。")
add_footer(slide, 7)

# ============================================================
# Slide 8: 竞品分析（碎片化）
# ============================================================
slide = add_content_slide("四、竞品动态分析", [
    "竞品A（数澜科技）：",
    "  - 完成C轮融资 3亿",
    "  - 发布「数据资产目录」新产品",
    "  - 挖走我方2名高级架构师",
    "",
    "竞品B（袋鼠云）：",
    "  - 与华为云达成战略合作",
    "  - 价格战：基础版免费",
    "",
    "竞品C（智领云）：",
    "  - 被某大厂收购（传闻未确认）",
    "  - 团队动荡，客户观望",
    "",
    "→ 详细竞品分析见附件《2024H1竞品跟踪报告》（另发）",
])
add_footer(slide, 8)

# ============================================================
# Slide 9: Q3计划
# ============================================================
slide = add_content_slide("五、Q3工作计划", [
    "重点方向：",
    "1. 秋季产品发布会（9月中旬，目标500人规模）",
    "2. 行业解决方案白皮书 x3（金融/制造/零售）",
    "3. 合作伙伴生态大会（联合10家ISV）",
    "4. 官网SEO优化（目标：自然流量+30%）",
    "5. 客户成功案例视频 x5",
    "",
    "里程碑：",
    "  7月：发布会筹备 + 白皮书启动",
    "  8月：合作伙伴招募 + 视频拍摄",
    "  9月：发布会执行 + 成果传播",
], notes="备注：发布会预算已获批80万，场地初步定在国家会议中心。需要产品部配合demo准备。")
add_footer(slide, 9)

# ============================================================
# Slide 10: 预算（数字为主）
# ============================================================
slide = add_content_slide("六、预算执行情况", [
    "Q2市场预算：总计 260万元",
    "",
    "项目          | 预算    | 实际    | 执行率",
    "────────────────────────────────────────",
    "线下活动      | 80万   | 76.5万  | 95.6%",
    "线上推广      | 60万   | 58.2万  | 97.0%",
    "品牌广告      | 50万   | 50万    | 100%",
    "内容制作      | 30万   | 27.8万  | 92.7%",
    "展会参展      | 25万   | 24.1万  | 96.4%",
    "其他          | 15万   | 12.3万  | 82.0%",
    "────────────────────────────────────────",
    "合计          | 260万  | 248.9万 | 95.7%",
    "",
    "Q3预算申请：310万（含发布会80万专项）",
])
add_footer(slide, 10)

# ============================================================
# Slide 11: 模板残留页（占位符未清理）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide.shapes.add_textbox(Cm(3), Cm(3), Cm(28), Cm(10))
tf = txBox.text_frame
placeholders = [
    "Click to edit title",
    "",
    "• Click to edit subtitle",
    "• 在此处添加要点",
    "• Add your key points here",
    "",
    "[LOGO PLACEHOLDER - 替换为公司logo]",
    "[CHART PLACEHOLDER - 插入数据图表]",
    "",
    "Template: XingChen_Corporate_2024_v3.pptx",
    "Designer: brand@xingchen-tech.com",
]
for i, line in enumerate(placeholders):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(14)
    p.font.color.rgb = GREY
add_footer(slide, 11)

# ============================================================
# Slide 12: 感谢页 + 联系方式噪声
# ============================================================
slide = add_title_slide(
    "感谢聆听",
    "THANK YOU"
)
txBox = slide.shapes.add_textbox(Cm(5), Cm(10), Cm(24), Cm(6))
tf = txBox.text_frame
thanks_noise = [
    "汇报人：林小燕 | 市场营销部总监",
    "邮箱：linxiaoyan@xingchen-tech.com",
    "电话：+86-10-8888-3333 | 内线：6001",
    "",
    "星辰科技有限公司 STAR CHEN TECHNOLOGY CO., LTD.",
    "北京市海淀区中关村科技园区创新大厦18层 | 邮编 100080",
    "www.xingchen-tech.com",
    "",
    "本文档为星辰科技内部资料，包含商业机密信息。",
    "未经书面授权，不得以任何形式复制、传播或向第三方披露。",
    "违者将依法追究法律责任。",
    "",
    "文档编号：MKT-2024-Q2-REPORT | 版本 V1.3 | 2024-07-05",
    "Printed: 2024-07-05 14:23:15 by linxiaoyan@xingchen-tech.com",
]
for i, line in enumerate(thanks_noise):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(9)
    p.font.color.rgb = GREY
add_footer(slide, 12)

# ============================================================
# 保存
# ============================================================
filepath = os.path.join(OUTPUT_DIR, "2024Q2市场营销汇报.pptx")
prs.save(filepath)
print(f"已生成: {filepath}")
print(f"文件大小: {os.path.getsize(filepath) / 1024:.1f} KB")
print(f"Slides: {len(prs.slides)}")
