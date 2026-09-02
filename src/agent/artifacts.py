"""Generate downloadable office artifacts from agent-authored Markdown."""

from __future__ import annotations

import json
import os
import re
import uuid
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Literal
from xml.sax.saxutils import escape

ArtifactFormat = Literal["docx", "pdf", "pptx"]

_MAX_CONTENT_LENGTH = 250_000
_MAX_SLIDES = 60
_INVALID_FILENAME = re.compile(r'[<>:"/\\|?*\x00-\x1f]')


@dataclass(frozen=True)
class Artifact:
    id: str
    format: ArtifactFormat
    title: str
    filename: str
    size: int
    download_url: str
    slidev_filename: str | None = None
    slidev_url: str | None = None


def artifacts_root() -> Path:
    return Path(os.getenv("ARTIFACTS_DIR", "artifacts")).resolve()


def _safe_stem(value: str) -> str:
    stem = _INVALID_FILENAME.sub("-", value).strip(" .-")
    stem = re.sub(r"\s+", " ", stem)
    return (stem[:80].rstrip() or "generated-document")


def _filename(title: str, requested: str | None, suffix: str) -> str:
    candidate = Path(requested or title).stem
    return f"{_safe_stem(candidate)}.{suffix}"


def _markdown_blocks(content: str) -> list[tuple[str, str]]:
    blocks: list[tuple[str, str]] = []
    paragraph: list[str] = []

    def flush() -> None:
        if paragraph:
            blocks.append(("paragraph", " ".join(paragraph).strip()))
            paragraph.clear()

    for raw_line in content.splitlines():
        line = raw_line.strip()
        if not line:
            flush()
            continue
        heading = re.match(r"^(#{1,6})\s+(.+)$", line)
        bullet = re.match(r"^[-*+]\s+(.+)$", line)
        numbered = re.match(r"^\d+[.)]\s+(.+)$", line)
        if heading:
            flush()
            blocks.append((f"heading{len(heading.group(1))}", heading.group(2)))
        elif bullet:
            flush()
            blocks.append(("bullet", bullet.group(1)))
        elif numbered:
            flush()
            blocks.append(("number", numbered.group(1)))
        else:
            paragraph.append(line)
    flush()
    return blocks


def _plain_markdown(text: str) -> str:
    text = re.sub(r"!\[([^]]*)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"\[([^]]+)\]\([^)]*\)", r"\1", text)
    return re.sub(r"[*_`~]", "", text).strip()


def _generate_docx(path: Path, title: str, content: str) -> None:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt

    document = Document()
    section = document.sections[0]
    section.top_margin = section.bottom_margin = Pt(56)
    title_paragraph = document.add_heading(_plain_markdown(title), level=0)
    title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

    for kind, value in _markdown_blocks(content):
        text = _plain_markdown(value)
        if not text:
            continue
        if kind.startswith("heading"):
            level = min(int(kind[-1]), 3)
            document.add_heading(text, level=level)
        elif kind == "bullet":
            document.add_paragraph(text, style="List Bullet")
        elif kind == "number":
            document.add_paragraph(text, style="List Number")
        else:
            document.add_paragraph(text)

    document.core_properties.title = title
    document.core_properties.author = "Team Knowledge Base Agent"
    document.save(path)


def _generate_pdf(path: Path, title: str, content: str) -> None:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    font_name = "STSong-Light"
    pdfmetrics.registerFont(UnicodeCIDFont(font_name))
    styles = getSampleStyleSheet()
    body = ParagraphStyle(
        "AgentBody",
        parent=styles["BodyText"],
        fontName=font_name,
        fontSize=10.5,
        leading=17,
        textColor=colors.HexColor("#263238"),
        spaceAfter=7,
    )
    heading = {
        level: ParagraphStyle(
            f"AgentHeading{level}",
            parent=body,
            fontSize={1: 18, 2: 15, 3: 12.5}.get(level, 11),
            leading={1: 24, 2: 21, 3: 18}.get(level, 17),
            textColor=colors.HexColor("#123C4A"),
            spaceBefore=12,
            spaceAfter=7,
        )
        for level in range(1, 7)
    }
    title_style = ParagraphStyle(
        "AgentTitle",
        parent=heading[1],
        fontSize=24,
        leading=32,
        alignment=TA_CENTER,
        spaceAfter=24,
    )
    document = SimpleDocTemplate(
        str(path),
        pagesize=A4,
        rightMargin=50,
        leftMargin=50,
        topMargin=54,
        bottomMargin=54,
        title=title,
        author="Team Knowledge Base Agent",
    )
    story = [Paragraph(escape(_plain_markdown(title)), title_style)]
    for kind, value in _markdown_blocks(content):
        text = escape(_plain_markdown(value))
        if not text:
            continue
        if kind.startswith("heading"):
            story.extend([Paragraph(text, heading[int(kind[-1])]), Spacer(1, 2)])
        elif kind in {"bullet", "number"}:
            marker = "&#8226;" if kind == "bullet" else "&#8212;"
            story.append(Paragraph(f"{marker}&nbsp;&nbsp;{text}", body))
        else:
            story.append(Paragraph(text, body))
    document.build(story)


def _slide_sections(title: str, content: str) -> list[tuple[str, list[str]]]:
    raw_sections = re.split(r"^\s*---\s*$", content, flags=re.MULTILINE)
    sections: list[tuple[str, list[str]]] = []
    for index, raw_section in enumerate(raw_sections):
        lines = [line.strip() for line in raw_section.splitlines() if line.strip()]
        if not lines:
            continue
        heading_index = next(
            (i for i, line in enumerate(lines) if re.match(r"^#{1,3}\s+", line)),
            None,
        )
        if heading_index is None:
            slide_title = title if index == 0 else f"{title} {index + 1}"
            body = lines
        else:
            slide_title = re.sub(r"^#{1,3}\s+", "", lines[heading_index])
            body = lines[:heading_index] + lines[heading_index + 1 :]
        sections.append(
            (
                _plain_markdown(slide_title),
                [_plain_markdown(re.sub(r"^[-*+]\s+", "", line)) for line in body],
            )
        )
    return (sections or [(title, [])])[:_MAX_SLIDES]


def _generate_pptx(path: Path, title: str, content: str) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    sections = _slide_sections(title, content)

    for index, (slide_title, lines) in enumerate(sections):
        layout = presentation.slide_layouts[0 if index == 0 and not lines else 1]
        slide = presentation.slides.add_slide(layout)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(247, 249, 248)
        title_shape = slide.shapes.title
        title_shape.text = slide_title
        title_frame = title_shape.text_frame
        for paragraph in title_frame.paragraphs:
            paragraph.font.name = "Microsoft YaHei"
            paragraph.font.size = Pt(30 if lines else 36)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(18, 60, 74)
            paragraph.alignment = PP_ALIGN.LEFT if lines else PP_ALIGN.CENTER

        if lines and len(slide.placeholders) > 1:
            frame = slide.placeholders[1].text_frame
            frame.clear()
            for line_index, line in enumerate(lines[:12]):
                paragraph = frame.paragraphs[0] if line_index == 0 else frame.add_paragraph()
                paragraph.text = line
                paragraph.level = 0
                paragraph.font.name = "Microsoft YaHei"
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = RGBColor(38, 50, 56)
                paragraph.space_after = Pt(8)

    presentation.core_properties.title = title
    presentation.core_properties.author = "Team Knowledge Base Agent"
    presentation.save(path)


def _slidev_source(title: str, content: str) -> str:
    cleaned = content.strip()
    if not cleaned.startswith("#"):
        cleaned = f"# {title}\n\n{cleaned}"
    return (
        "---\n"
        "theme: default\n"
        f"title: {json.dumps(title, ensure_ascii=False)}\n"
        "download: true\n"
        "transition: slide-left\n"
        "---\n\n"
        f"{cleaned}\n"
    )


def generate_artifact(
    *,
    format: ArtifactFormat,
    title: str,
    content: str,
    file_name: str | None = None,
) -> Artifact:
    title = title.strip()
    content = content.strip()
    if not title:
        raise ValueError("文档标题不能为空")
    if not content:
        raise ValueError("文档内容不能为空")
    if len(content) > _MAX_CONTENT_LENGTH:
        raise ValueError(f"文档内容不能超过 {_MAX_CONTENT_LENGTH} 个字符")
    if format not in {"docx", "pdf", "pptx"}:
        raise ValueError(f"不支持的文档格式: {format}")

    artifact_id = str(uuid.uuid4())
    directory = artifacts_root() / artifact_id
    directory.mkdir(parents=True, exist_ok=False)
    filename = _filename(title, file_name, format)
    path = directory / filename

    try:
        if format == "docx":
            _generate_docx(path, title, content)
        elif format == "pdf":
            _generate_pdf(path, title, content)
        else:
            _generate_pptx(path, title, content)

        slidev_filename = None
        slidev_url = None
        if format == "pptx":
            slidev_filename = f"{Path(filename).stem}.slides.md"
            (directory / slidev_filename).write_text(
                _slidev_source(title, content), encoding="utf-8"
            )
            slidev_url = f"/api/artifacts/{artifact_id}/slidev"

        artifact = Artifact(
            id=artifact_id,
            format=format,
            title=title,
            filename=filename,
            size=path.stat().st_size,
            download_url=f"/api/artifacts/{artifact_id}/download",
            slidev_filename=slidev_filename,
            slidev_url=slidev_url,
        )
        (directory / "metadata.json").write_text(
            json.dumps(asdict(artifact), ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        return artifact
    except Exception:
        for child in directory.iterdir():
            child.unlink(missing_ok=True)
        directory.rmdir()
        raise


def resolve_artifact(artifact_id: str, *, slidev: bool = False) -> tuple[Path, Artifact]:
    try:
        normalized_id = str(uuid.UUID(artifact_id))
    except ValueError as exc:
        raise FileNotFoundError("artifact not found") from exc
    directory = artifacts_root() / normalized_id
    metadata_path = directory / "metadata.json"
    try:
        artifact = Artifact(**json.loads(metadata_path.read_text(encoding="utf-8")))
    except (OSError, TypeError, ValueError, json.JSONDecodeError) as exc:
        raise FileNotFoundError("artifact not found") from exc
    filename = artifact.slidev_filename if slidev else artifact.filename
    if not filename:
        raise FileNotFoundError("artifact variant not found")
    path = directory / filename
    if not path.is_file() or path.parent != directory:
        raise FileNotFoundError("artifact not found")
    return path, artifact
